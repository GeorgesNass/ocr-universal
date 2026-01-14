'''
__author__ = "Georges Nassopoulos"
__copyright__ = None
__version__ = "1.0.0"
__email__ = "georges.nassopoulos@gmail.com"
__status__ = "Prod"
__desc__ = "Extract text from PPT and PPTX files, with optional detailed font attribute extraction."
'''

from pathlib import Path
import os
import subprocess
from pptx import Presentation
from src.utils import get_logger, log_execution_time_and_path
from src.utils.constants import (
    INPUT_DIR,
    CONVERTED_DIR,
    OUTPUT_DIR,
    PATH_LIBRE_OFFICE,
    USE_DETAILED_PPTX_EXTRACTION
)
from src.utils.ocr_utils import is_allowed_file

## ============================================================
## LOGGER
## ============================================================
logger = get_logger("pptx_ppt_to_text")

## ============================================================
## FUNCTION: Extract text (simple)
## ============================================================
def extract_text_simple(pptx_path: Path) -> str:
    """
        Extract only visible text from all slides without formatting details

        Args:
            pptx_path (Path): Path to the .pptx file

        Returns:
            str: Extracted plain text from all slides
    """
    
    presentation = Presentation(pptx_path)
    text_runs = []

    ## Loop through all slides and shapes
    for slide in presentation.slides:
        for shape in slide.shapes:
        
            ## Skip shapes without text
            if not shape.has_text_frame:
                continue
                
            ## Extract plain text if not empty
            if shape.text.strip():
                text_runs.append(shape.text.strip())

    logger.info(f"Total slides processed: {len(presentation.slides)}, total text length: {sum(len(t) for t in text_runs)} chars")

    return "\n\n".join(text_runs)

## ============================================================
## FUNCTION: Extract text (detailed, with formatting)
## ============================================================
def extract_text_detailed(pptx_path: Path) -> str:
    """
        Extract text from all slides while preserving font attributes
        (bold, italic, color, font name, size)

        Args:
            pptx_path (Path): Path to the .pptx file

        Returns:
            str: Extracted text content from all slides
    """
    
    presentation = Presentation(pptx_path)
    text_runs = []

    ## Iterate over each slide
    for slide in presentation.slides:
    
        ## Iterate over each shape in the slide
        for shape in slide.shapes:
        
            ## Skip if no text
            if not shape.has_text_frame:
                continue

            ## For each paragraph in the text frame
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font = run.font
                    text_content = run.text

                    ## Collect optional font attributes
                    try:
                        color_attr = font.color.rgb if font.color and font.color.type else None
                    except Exception:
                        color_attr = None

                    ## Log style information (optional)
                    logger.debug(
                        f"Run: '{text_content[:30]}...' "
                        f"[bold={font.bold}, italic={font.italic}, "
                        f"font={font.name}, size={font.size}, color={color_attr}]"
                    )

                    text_runs.append(text_content)

    logger.info(f"Extracted detailed text from {len(presentation.slides)} slides, total length: {sum(len(t) for t in text_runs)} chars")

    return "\n\n".join(text_runs)

## ============================================================
## FUNCTION: Extract PPTX text (with selection)
## ============================================================
@log_execution_time_and_path
def get_text_from_pptx(pptx_path: Path) -> str:
    """
        Main function to extract text from a PowerPoint (.pptx) file
        Chooses between simple or detailed extraction mode depending on configuration

        Args:
            pptx_path (Path): Path to the PowerPoint file

        Returns:
            str: Extracted text
    """
    
    try:
        ## Choose extraction mode based on global flag
        if USE_DETAILED_PPTX_EXTRACTION:
            logger.info("Using detailed PPTX text extraction (with font attributes)")
            return extract_text_detailed(pptx_path)
        else:
            logger.info("Using simple PPTX text extraction (plain text only)")
            return extract_text_simple(pptx_path)

    except Exception as e:
        logger.exception(f"Failed to extract text from {pptx_path}: {e}")
        return ""

## ============================================================
## FUNCTION: Convert PPT to PPTX (LibreOffice)
## ============================================================
@log_execution_time_and_path
def convert_ppt_to_pptx(src_path: str) -> Path:
    """
        Convert an old PowerPoint (.ppt) file into .pptx format using LibreOffice

        Args:
            src_path (Path): Path to the file

        Returns:
            Path: Path to the converted .pptx file
    """
    
    ## Keep homogeneous exists check
    if not os.path.exists(str(src_path)):
        logger.error(f"Source file not found: {src_path}")
        return None

    ## Normalize to Path only for suffix/stem building
    ppt_path = Path(src_path)

    try:
        ## Build output file path
        output_path = CONVERTED_DIR / f"{ppt_path.stem}.pptx"

        ## Select LibreOffice executable (Windows/Linux)
        ## PATH_LIBRE_OFFICE should be:
        ## - Windows: C:\\Program Files\\LibreOffice\\program\\soffice.exe
        ## - Linux: soffice (or /usr/bin/libreoffice)
        libreoffice_exec = PATH_LIBRE_OFFICE

        ## Run LibreOffice conversion (secure subprocess)
        result = subprocess.run(
            [
                str(libreoffice_exec),
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                str(CONVERTED_DIR),
                str(ppt_path)
            ],
            capture_output=True,
            text=True
        )

        ## Check return code
        if result.returncode != 0:
            logger.error(f"LibreOffice conversion failed: {result.stderr}")
            return None

        ## Verify conversion result (keep os.path.exists for homogeneity)
        if not os.path.exists(str(output_path)):
            logger.error(f"Conversion failed: {output_path} not found.")
            return None

        logger.info(f"Conversion successful: {output_path}")
        return output_path

    except Exception as e:
        logger.exception(f"Error converting {ppt_path} to pptx: {e}")
        return None

## ============================================================
## FUNCTION: Process PPT/PPTX file end-to-end
## ============================================================
@log_execution_time_and_path
def process_presentation(src_path: str) -> None:
    """
        Full pipeline for processing PowerPoint files:
            - Detects if it's .ppt or .pptx
            - Converts if needed
            - Extracts text
            - Saves to output folder

        Args:
            src_path (str): File path complet to process
    """
    
    ## Keep homogeneous exists check
    if not os.path.exists(str(src_path)):
        logger.error(f"Input file not found: {src_path}")
        return

    ## Normalize to Path only for suffix/stem
    src_path = Path(src_path)

    ## If PPT, convert to PPTX first
    if src_path.suffix.lower() == ".ppt":
        logger.info(f"Converting .ppt to .pptx: {src_path}")
        converted = convert_ppt_to_pptx(str(src_path))
        if not converted:
            return
        src_path = Path(converted)

    ## Extract text
    text = get_text_from_pptx(src_path)

    ## Build output path
    output_path = OUTPUT_DIR / f"{src_path.stem}.txt"

    ## Skip empty outputs
    if not text.strip():
        logger.warning(f"No text extracted from {src_path.name}, skipping save.")
        return

    ## Save text
    try:
        output_path.write_text(text, encoding="utf-8")
        logger.info(f"Text successfully saved to: {output_path}")
    except Exception as e:
        logger.exception(f"Error saving extracted text: {e}")
        
## ============================================================
## MAIN ENTRY POINT
## ============================================================
if __name__ == "__main__":
    """
        Example usage when running this script directly
        Processes all PPT and PPTX files found in the input directory
    """

    logger.info("Starting manual execution: PowerPoint text extraction")

    ## Defensive: ensure INPUT_DIR exists
    if not INPUT_DIR.exists():
        logger.error(f"INPUT_DIR does not exist: {INPUT_DIR}")
        raise SystemExit(1)

    ## Process only PPT/PPTX
    for file in INPUT_DIR.glob("*"):
        if not file.is_file():
            continue

        if file.suffix.lower() in [".ppt", ".pptx"]:
            logger.info(f"Processing file: {file}")
            try:
                ## Ensure Path is passed (process_presentation must accept Path)
                process_presentation(file)
            except Exception as e:
                logger.exception(f"Failed processing file {file}: {e}")

    logger.info("Finished manual processing of PowerPoint files.")
